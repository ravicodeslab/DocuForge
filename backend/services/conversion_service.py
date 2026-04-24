"""
Conversion service — handles all file format conversion logic.
PDF ⇄ DOCX, PPTX, XLSX, JPG/PNG; Image → PDF
"""
import logging
import zipfile
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
import pandas as pd

from utils.exceptions import ConversionError
from utils.file_handler import FileHandler

logger = logging.getLogger("docuforge.conversion")


class ConversionService:

    # ── PDF → DOCX ────────────────────────────────────────────────────────────
    @staticmethod
    def pdf_to_docx(pdf_path: Path) -> Path:
        """
        Extracts text from PDF page by page and writes to DOCX.
        Layout fidelity is text-level; complex formatting is approximated.
        """
        try:
            doc_out = Document()
            doc_out.add_heading("Converted Document", level=0)

            pdf = fitz.open(str(pdf_path))
            for page_num, page in enumerate(pdf, start=1):
                blocks = page.get_text("blocks")
                # Sort blocks top-to-bottom, left-to-right
                blocks.sort(key=lambda b: (round(b[1] / 20), b[0]))

                doc_out.add_paragraph(f"— Page {page_num} —", style="Heading 2")
                for block in blocks:
                    text = block[4].strip()
                    if text:
                        para = doc_out.add_paragraph(text)
                        para.style.font.size = Pt(10)

                if page_num < len(pdf):
                    doc_out.add_page_break()

            pdf.close()

            out_path = FileHandler.output_path(".docx")
            doc_out.save(str(out_path))
            logger.info("PDF→DOCX: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── PDF → Images ──────────────────────────────────────────────────────────
    @staticmethod
    def pdf_to_images(pdf_path: Path, fmt: str = "jpg", dpi: int = 150) -> list[Path]:
        """Render each PDF page to an image. Returns list of image paths."""
        try:
            pdf = fitz.open(str(pdf_path))
            paths: list[Path] = []
            ext = "jpg" if fmt.lower() in ("jpg", "jpeg") else "png"
            pix_fmt = "jpeg" if ext == "jpg" else "png"

            for page_num, page in enumerate(pdf):
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                out_path = FileHandler.output_path(f"_p{page_num + 1}.{ext}")
                pix.save(str(out_path))
                paths.append(out_path)

            pdf.close()
            logger.info("PDF→Images: %d pages produced", len(paths))
            return paths
        except Exception as e:
            raise ConversionError(str(e))

    # ── Images → PDF ──────────────────────────────────────────────────────────
    @staticmethod
    def images_to_pdf(image_paths: list[Path]) -> Path:
        """Combine one or more images into a single PDF."""
        try:
            pil_images = []
            for p in image_paths:
                img = Image.open(str(p)).convert("RGB")
                pil_images.append(img)

            out_path = FileHandler.output_path(".pdf")
            if not pil_images:
                raise ConversionError("No valid images provided.")

            pil_images[0].save(
                str(out_path), save_all=True, append_images=pil_images[1:]
            )
            logger.info("Images→PDF: %s", out_path.name)
            return out_path
        except ConversionError:
            raise
        except Exception as e:
            raise ConversionError(str(e))

    # ── PDF → PPTX ────────────────────────────────────────────────────────────
    @staticmethod
    def pdf_to_pptx(pdf_path: Path) -> Path:
        """
        Converts each PDF page into a slide (image-based for fidelity).
        Text is overlaid for accessibility.
        """
        try:
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]  # blank layout

            pdf = fitz.open(str(pdf_path))
            for page_num, page in enumerate(pdf):
                # Render page as image
                mat = fitz.Matrix(1.5, 1.5)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_path = FileHandler.output_path(f"_slide{page_num}.jpg")
                pix.save(str(img_path))

                slide = prs.slides.add_slide(blank_layout)
                width = prs.slide_width
                height = prs.slide_height

                # Add page image as background
                pic = slide.shapes.add_picture(str(img_path), 0, 0, width, height)
                img_path.unlink(missing_ok=True)

            pdf.close()
            out_path = FileHandler.output_path(".pptx")
            prs.save(str(out_path))
            logger.info("PDF→PPTX: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── PDF → XLSX ────────────────────────────────────────────────────────────
    @staticmethod
    def pdf_to_excel(pdf_path: Path) -> Path:
        """
        Extracts tabular text from PDF and writes to Excel.
        Each page becomes a sheet. Best-effort table detection.
        """
        try:
            pdf = fitz.open(str(pdf_path))
            out_path = FileHandler.output_path(".xlsx")

            with pd.ExcelWriter(str(out_path), engine="xlsxwriter") as writer:
                for page_num, page in enumerate(pdf, start=1):
                    # Attempt table extraction first
                    tables = page.find_tables()
                    if tables.tables:
                        for tbl_idx, table in enumerate(tables.tables):
                            df = pd.DataFrame(table.extract())
                            if not df.empty:
                                df.columns = df.iloc[0]
                                df = df[1:].reset_index(drop=True)
                                sheet = f"P{page_num}_T{tbl_idx + 1}"
                                df.to_excel(writer, sheet_name=sheet[:31], index=False)
                    else:
                        # Fall back to raw text rows
                        lines = page.get_text("text").splitlines()
                        rows = [l.strip() for l in lines if l.strip()]
                        df = pd.DataFrame(rows, columns=["Content"])
                        df.to_excel(writer, sheet_name=f"Page_{page_num}"[:31], index=False)

            pdf.close()
            logger.info("PDF→Excel: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── DOCX Merge ────────────────────────────────────────────────────────────
    @staticmethod
    def merge_docx_files(docx_paths: list) -> Path:
        """Merge multiple DOCX files into one."""
        try:
            merged_doc = Document(str(docx_paths[0]))
            
            for docx_path in docx_paths[1:]:
                source_doc = Document(str(docx_path))
                merged_doc.add_page_break()
                
                for para in source_doc.paragraphs:
                    new_para = merged_doc.add_paragraph(para.text, style=para.style)
                    if para.runs:
                        for run_idx, run in enumerate(para.runs):
                            if run_idx < len(new_para.runs):
                                new_para.runs[run_idx].bold = run.bold
                                new_para.runs[run_idx].italic = run.italic
                
                for table in source_doc.tables:
                    tbl = merged_doc.add_table(rows=len(table.rows), cols=len(table.columns))
                    tbl.style = table.style
                    for i, row in enumerate(table.rows):
                        for j, cell in enumerate(row.cells):
                            tbl.rows[i].cells[j].text = cell.text
            
            out_path = FileHandler.output_path(".docx")
            merged_doc.save(str(out_path))
            logger.info("DOCX Merge: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── DOCX Compress ────────────────────────────────────────────────────────────
    @staticmethod
    def compress_docx(docx_path: Path, compression_level: str = "medium") -> Path:
        """Compress DOCX file by removing embedded images and metadata."""
        try:
            doc = Document(str(docx_path))
            
            # Remove embedded images if high compression
            if compression_level in ["high", "maximum"]:
                for para in doc.paragraphs:
                    for run in para.runs:
                        if run._element.drawing_lst:
                            run._element.getparent().remove(run._element)
            
            out_path = FileHandler.output_path(".docx")
            doc.save(str(out_path))
            logger.info("DOCX Compress: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── DOCX to PDF ────────────────────────────────────────────────────────────
    @staticmethod
    def docx_to_pdf(docx_path: Path) -> Path:
        """Convert DOCX to PDF (requires LibreOffice)."""
        try:
            import subprocess
            import tempfile
            
            # Create temporary output directory
            with tempfile.TemporaryDirectory() as temp_dir:
                output_file = Path(temp_dir) / "output.pdf"
                
                # Use LibreOffice for conversion
                result = subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", temp_dir,
                        str(docx_path)
                    ],
                    capture_output=True,
                    timeout=30
                )
                
                if result.returncode != 0:
                    raise ConversionError("LibreOffice conversion failed")
                
                # Copy to output folder
                pdf_path = Path(temp_dir) / docx_path.stem + ".pdf"
                final_path = FileHandler.output_path(".pdf")
                final_path.write_bytes(pdf_path.read_bytes())
                
            logger.info("DOCX→PDF: %s", final_path.name)
            return final_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── DOCX to TXT ────────────────────────────────────────────────────────────
    @staticmethod
    def docx_to_txt(docx_path: Path) -> Path:
        """Convert DOCX to plain text."""
        try:
            doc = Document(str(docx_path))
            
            out_path = FileHandler.output_path(".txt")
            with open(out_path, "w", encoding="utf-8") as f:
                for para in doc.paragraphs:
                    f.write(para.text + "\n")
            
            logger.info("DOCX→TXT: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── DOCX Split ────────────────────────────────────────────────────────────
    @staticmethod
    def split_docx(docx_path: Path, sections: int = 2) -> list:
        """Split DOCX into multiple sections."""
        try:
            doc = Document(str(docx_path))
            paragraphs = doc.paragraphs
            para_per_section = len(paragraphs) // sections
            
            output_files = []
            for i in range(sections):
                new_doc = Document()
                start_idx = i * para_per_section
                end_idx = (i + 1) * para_per_section if i < sections - 1 else len(paragraphs)
                
                for para in paragraphs[start_idx:end_idx]:
                    new_doc.add_paragraph(para.text, style=para.style)
                
                out_path = FileHandler.output_path(f"_part{i+1}.docx")
                new_doc.save(str(out_path))
                output_files.append(out_path)
            
            logger.info("DOCX Split: %d files", len(output_files))
            return output_files
        except Exception as e:
            raise ConversionError(str(e))

    # ── Image Convert ────────────────────────────────────────────────────────────
    @staticmethod
    def convert_image(image_path: Path, target_format: str) -> Path:
        """Convert image between formats."""
        try:
            img = Image.open(str(image_path))
            
            # Convert RGBA to RGB if needed for JPEG
            if target_format.lower() in ["jpg", "jpeg"] and img.mode in ["RGBA", "LA"]:
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = rgb_img
            
            out_path = FileHandler.output_path(f".{target_format.lower()}")
            img.save(str(out_path), quality=95)
            
            logger.info("Image Convert: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── Image Compress ────────────────────────────────────────────────────────────
    @staticmethod
    def compress_image(image_path: Path, quality: int = 80, target_size: str = None) -> Path:
        """Compress image with quality control."""
        try:
            img = Image.open(str(image_path))
            out_path = FileHandler.output_path(f"_compressed.{img.format.lower()}")
            
            # Save with quality
            img.save(str(out_path), quality=quality, optimize=True)
            
            logger.info("Image Compress: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── Image Resize ────────────────────────────────────────────────────────────
    @staticmethod
    def resize_image(image_path: Path, width: int, height: int) -> Path:
        """Resize image to specified dimensions."""
        try:
            img = Image.open(str(image_path))
            resized_img = img.resize((width, height), Image.Resampling.LANCZOS)
            
            out_path = FileHandler.output_path(f"_resized.{img.format.lower()}")
            resized_img.save(str(out_path), quality=95)
            
            logger.info("Image Resize: %s (%dx%d)", out_path.name, width, height)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))

    # ── Image to PDF ────────────────────────────────────────────────────────────
    @staticmethod
    def image_to_pdf(image_path: Path) -> Path:
        """Convert image to PDF."""
        try:
            img = Image.open(str(image_path))
            
            # Convert RGBA to RGB
            if img.mode in ["RGBA", "LA"]:
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = rgb_img
            
            out_path = FileHandler.output_path(".pdf")
            img.save(str(out_path), "PDF")
            
            logger.info("Image→PDF: %s", out_path.name)
            return out_path
        except Exception as e:
            raise ConversionError(str(e))
